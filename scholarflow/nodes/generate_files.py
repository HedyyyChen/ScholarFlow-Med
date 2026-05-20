"""生成文件节点 - Excel 和 PPT"""

import os
from typing import Dict, Any, Optional
from scholarflow.research_state import ResearchState


def _determine_author_position(authors: list, expert_name: str) -> str:
    """判断作者排位 (first / corresponding / other)"""
    if not authors or not expert_name:
        return "other"

    name_lower = expert_name.lower()
    name_parts = name_lower.split()

    # 检查是否是第一作者
    first_author = (authors[0] or "").lower()
    if any(part in first_author for part in name_parts if len(part) > 1):
        return "first"

    # 检查是否是通讯作者（通常是最后一位）
    if len(authors) >= 2:
        last_author = (authors[-1] or "").lower()
        if any(part in last_author for part in name_parts if len(part) > 1):
            return "corresponding"

    # 检查是否在任何位置
    for author in authors:
        author_lower = (author or "").lower()
        if any(part in author_lower for part in name_parts if len(part) > 1):
            return "other"

    return "other"


def _generate_summary_with_llm(title: str, abstract: str = "") -> str:
    """使用 DeepSeek LLM 生成中文核心内容简述"""
    try:
        from scholarflow.tools.llm_client import get_client
        client = get_client()
        if client:
            return client.generate_abstract_summary(title, abstract)
    except Exception:
        pass
    return ""


def _generate_abstract_summary_fallback(title: str, abstract: str = "") -> str:
    """回退方案：基于标题和摘要生成简单中文简述"""
    title_lower = (title or "").lower()
    abstract_lower = (abstract or "").lower()
    text = f"{title} {abstract_lower}"

    # 关键词映射
    keywords = {
        "骨髓瘤": ["multiple myeloma", "myeloma"],
        "淀粉样变性": ["amyloidosis"],
        "白血病": ["leukemia"],
        "淋巴瘤": ["lymphoma"],
        "CAR-T": ["car-t", "car t", "chimeric antigen receptor"],
        "造血干细胞移植": ["transplant", "stem cell"],
        "靶向治疗": ["targeted therapy"],
        "免疫治疗": ["immunotherapy"],
    }
    detected = []
    for cn, en_list in keywords.items():
        if any(kw in text for kw in en_list):
            detected.append(cn)
    if detected:
        return f"本研究聚焦{'、'.join(detected)}领域，探讨" + title[:50] + "..."
    return title[:50] + "的研究"


def _get_impact_factor(journal: str) -> Optional[str]:
    """根据期刊名获取 2025 影响因子（硬编码常见期刊）"""
    if not journal:
        return None

    journal_lower = journal.lower()

    if_factors = {
        "blood": "21.0",
        "new england journal of medicine": "96.2",
        "lancet": "168.9",
        "nature": "50.5",
        "science": "44.7",
        "journal of clinical oncology": "42.1",
        "nature medicine": "58.7",
        "cell": "45.5",
        "leukemia": "11.4",
        "clinical cancer research": "11.5",
        "american journal of hematology": "10.6",
        "blood cancer journal": "9.5",
        "haematologica": "9.1",
        "british journal of haematology": "5.5",
        "bone marrow transplantation": "4.8",
        "transplantation": "4.0",
        "biology of blood and marrow transplantation": "4.2",
        "frontiers in immunology": "5.7",
        "journal of hematology & oncology": "17.4",
        "cancers": "4.5",
        "plos one": "2.9",
        "scientific reports": "3.8",
        "european journal of haematology": "3.1",
        "annals of hematology": "3.5",
        "hematological oncology": "3.1",
        "current opinion in hematology": "3.9",
        "best practice & research clinical haematology": "2.5",
        "clinical lymphoma myeloma and leukemia": "2.5",
        "journal of clinical medicine": "3.0",
        "oncotarget": "3.0",
        "medicine": "1.6",
        "bmc cancer": "3.8",
        "journal of cancer research and clinical oncology": "4.0",
        "leukemia research": "2.3",
        "international journal of molecular sciences": "5.6",
        "npj precision oncology": "7.8",
        "nature reviews clinical oncology": "81.1",
        "nature reviews drug discovery": "112.0",
        "signal transduction and targeted therapy": "39.3",
        "cell discovery": "13.0",
        "advanced science": "17.7",
        "nature communications": "14.7",
        "proceedings of the national academy of sciences": "9.4",
        "elife": "7.7",
        "jama": "56.7",
        "jama oncology": "27.4",
        "the lancet oncology": "52.9",
        "the lancet haematology": "18.9",
        "blood advances": "7.4",
        "ebiomedicine": "8.7",
        "frontiers in oncology": "3.5",
        "clinical and experimental medicine": "4.7",
        "experimental hematology": "3.4",
        "medical oncology": "2.8",
        "oncology letters": "2.5",
        "molecular medicine reports": "2.9",
        "experimental and therapeutic medicine": "2.4",
        "oncology reports": "3.8",
        "international journal of hematology": "2.8",
        "journal of clinical oncology": "42.1",
    }

    for key, value in if_factors.items():
        if key in journal_lower:
            return value

    return "1.0"


def generate_excel_1(state: ResearchState) -> Dict[str, Any]:
    """生成 1.xlsx - 专家论文Excel文件"""
    papers = state.get("filtered_expert_papers", [])
    expert_name = state.get("expert_name", "")
    output_dir = state.get("output_dir", "output")

    if not papers:
        return {"error": "没有论文数据，无法生成Excel"}

    os.makedirs(output_dir, exist_ok=True)
    excel_path = os.path.join(output_dir, "1.xlsx")

    try:
        import openpyxl
        from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
        from openpyxl.utils import get_column_letter

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "专家论文"

        # 表头样式
        header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
        header_font = Font(color="FFFFFF", bold=True, size=11)
        header_alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        thin_border = Border(
            left=Side(style="thin"), right=Side(style="thin"),
            top=Side(style="thin"), bottom=Side(style="thin")
        )

        # 表头：按照 workbuddy 格式
        headers = ["序号", "论文标题", "期刊", "2025 IF", "发表年份", "作者排位", "核心研究内容简述", "DOI"]
        col_widths = [8, 60, 30, 12, 12, 15, 50, 30]

        for col, header in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col, value=header)
            cell.fill = header_fill
            cell.font = header_font
            cell.alignment = header_alignment
            cell.border = thin_border
            ws.column_dimensions[get_column_letter(col)].width = col_widths[col - 1]

        for idx, paper in enumerate(papers, 1):
            row = idx + 1

            title = paper.get("title", "")
            journal = paper.get("journal", "")
            year = paper.get("year", "")

            # 获取 IF（不能为空）
            impact_factor = paper.get("impact_factor")
            if not impact_factor:
                impact_factor = _get_impact_factor(journal)
            if not impact_factor:
                impact_factor = "1.0"

            # 判断作者排位（不能为空）
            author_position = paper.get("author_position")
            if not author_position:
                author_position = _determine_author_position(paper.get("authors", []), expert_name)

            # 生成中文核心研究内容简述（不能为空）
            summary = _generate_summary_with_llm(title, paper.get("abstract", ""))
            if not summary or summary == title[:50]:
                summary = _generate_abstract_summary_fallback(title, paper.get("abstract", ""))

            # DOI 加链接格式
            doi = paper.get("doi") or ""
            doi_display = doi
            if doi and not doi.startswith("http"):
                doi_display = f"https://doi.org/{doi}"

            ws.cell(row=row, column=1, value=idx)
            ws.cell(row=row, column=2, value=title)
            ws.cell(row=row, column=3, value=journal)
            if_cell = ws.cell(row=row, column=4, value=impact_factor)
            if_cell.alignment = Alignment(horizontal="center", vertical="top")
            year_cell = ws.cell(row=row, column=5, value=year if year else "")
            year_cell.alignment = Alignment(horizontal="center", vertical="top")
            ws.cell(row=row, column=6, value=author_position)
            ws.cell(row=row, column=7, value=summary[:80])
            ws.cell(row=row, column=8, value=doi_display)

            for col in range(1, 9):
                cell = ws.cell(row=row, column=col)
                cell.border = thin_border
                cell.alignment = Alignment(wrap_text=True, vertical="top")

        # 冻结首行
        ws.freeze_panes = "A2"

        wb.save(excel_path)
        print(f"  -> 已生成 {excel_path}")
        return {"excel_1_path": excel_path, "current_step": "generate_excel_1_completed"}

    except Exception as e:
        return {"error": f"生成Excel 1失败: {str(e)}"}


def generate_excel_2(state: ResearchState) -> Dict[str, Any]:
    """生成 2.xlsx - AI相关论文Excel文件"""
    papers = state.get("ai_papers", [])
    output_dir = state.get("output_dir", "output")

    if not papers:
        return {"error": "没有AI论文数据，无法生成Excel"}

    os.makedirs(output_dir, exist_ok=True)
    excel_path = os.path.join(output_dir, "2.xlsx")

    try:
        import openpyxl
        from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
        from openpyxl.utils import get_column_letter

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "AI相关论文"

        header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
        header_font = Font(color="FFFFFF", bold=True, size=11)
        header_alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        thin_border = Border(
            left=Side(style="thin"), right=Side(style="thin"),
            top=Side(style="thin"), bottom=Side(style="thin")
        )

        headers = ["序号", "论文标题", "期刊", "2025 IF", "发表年份", "作者排位", "核心研究内容简述", "DOI", "研究领域", "AI技术"]
        col_widths = [8, 60, 30, 12, 12, 15, 50, 30, 20, 20]

        for col, header in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col, value=header)
            cell.fill = header_fill
            cell.font = header_font
            cell.alignment = header_alignment
            cell.border = thin_border
            ws.column_dimensions[get_column_letter(col)].width = col_widths[col - 1]

        for idx, paper in enumerate(papers, 1):
            row = idx + 1

            journal = paper.get("journal", "")
            impact_factor = _get_impact_factor(journal) or "1.0"

            title = paper.get("title", "")
            summary = _generate_summary_with_llm(title, paper.get("abstract", ""))
            if not summary:
                summary = _generate_abstract_summary_fallback(title, paper.get("abstract", ""))

            doi = paper.get("doi") or ""
            doi_display = doi
            if doi and not doi.startswith("http"):
                doi_display = f"https://doi.org/{doi}"

            ws.cell(row=row, column=1, value=idx)
            ws.cell(row=row, column=2, value=title)
            ws.cell(row=row, column=3, value=journal)
            if_cell = ws.cell(row=row, column=4, value=impact_factor)
            if_cell.alignment = Alignment(horizontal="center", vertical="top")
            year_cell = ws.cell(row=row, column=5, value=paper.get("year", ""))
            year_cell.alignment = Alignment(horizontal="center", vertical="top")
            ws.cell(row=row, column=6, value="other")
            ws.cell(row=row, column=7, value=summary[:80])
            ws.cell(row=row, column=8, value=doi_display)
            ws.cell(row=row, column=9, value=paper.get("source", ""))
            ws.cell(row=row, column=10, value="AI/ML")

            for col in range(1, 11):
                cell = ws.cell(row=row, column=col)
                cell.border = thin_border
                cell.alignment = Alignment(wrap_text=True, vertical="top")

        ws.freeze_panes = "A2"
        wb.save(excel_path)
        print(f"  -> 已生成 {excel_path}")
        return {"excel_2_path": excel_path, "current_step": "generate_excel_2_completed"}

    except Exception as e:
        return {"error": f"生成Excel 2失败: {str(e)}"}


def generate_ppt(state: ResearchState) -> Dict[str, Any]:
    """生成 out.pptx - 医工交叉合作建议幻灯片"""
    expert_name = state.get("expert_name", "")
    institution = state.get("institution", "")
    research_areas = state.get("research_areas", [])
    research_areas_detail = state.get("research_areas_detail", [])
    expert_papers = state.get("filtered_expert_papers", [])
    ai_papers = state.get("ai_papers", [])
    topics = state.get("collaboration_topics", [])
    output_dir = state.get("output_dir", "output")

    os.makedirs(output_dir, exist_ok=True)
    ppt_path = os.path.join(output_dir, "out.pptx")

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # 颜色定义
        DARK_BG = RGBColor(0x1F, 0x29, 0x37)
        ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)
        WHITE = RGBColor(0xFF, 0xFF, 0xFF)
        LIGHT_GRAY = RGBColor(0xF3, 0xF4, 0xF6)

        def add_dark_slide():
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = DARK_BG
            return slide

        def add_textbox(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
            txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.bold = bold
            p.alignment = alignment
            return txBox

        # ---- Slide 1: Cover ----
        slide = add_dark_slide()
        add_textbox(slide, 1, 1.5, 11, 1.5, "医工交叉合作建议书", 44, WHITE, True, PP_ALIGN.CENTER)
        add_textbox(slide, 1, 3.2, 11, 0.8, f"专家: {expert_name}  |  机构: {institution}", 24, ACCENT_BLUE, False, PP_ALIGN.CENTER)
        add_textbox(slide, 1, 4.5, 11, 0.5, "基于文献分析的医工交叉合作方向建议", 18, LIGHT_GRAY, False, PP_ALIGN.CENTER)

        # ---- Slide 2: Academic Profile ----
        slide = add_dark_slide()
        add_textbox(slide, 0.8, 0.3, 11, 0.8, "专家学术画像", 32, WHITE, True)

        # 显示详细研究领域
        if research_areas_detail:
            detail_lines = [
                f"  - {a.get('领域中文名', a.get('领域英文名', ''))}: {a.get('论文数量', 0)}篇 ({a.get('占比', '')})"
                for a in research_areas_detail[:3]
            ]
            research_text = "主要研究领域:\n" + "\n".join(detail_lines)
        else:
            research_text = "主要研究领域:\n" + "\n".join([f"  - {area}" for area in research_areas])
        add_textbox(slide, 0.8, 1.3, 11, 2.5, research_text, 18, LIGHT_GRAY)
        stats_text = f"论文数量: {len(expert_papers)}篇 (近5年)"
        add_textbox(slide, 0.8, 4.0, 11, 1.5, stats_text, 18, LIGHT_GRAY)

        # ---- Slide 3: Paper Statistics ----
        slide = add_dark_slide()
        add_textbox(slide, 0.8, 0.3, 11, 0.8, "论文统计", 32, WHITE, True)
        year_count = {}
        for p in expert_papers:
            year = p.get("year")
            if year:
                year_count[year] = year_count.get(year, 0) + 1
        stats = "\n".join([f"{year}年: {count}篇" for year, count in sorted(year_count.items())])
        add_textbox(slide, 0.8, 1.5, 11, 3, f"论文年份分布:\n{stats}", 18, LIGHT_GRAY)

        # ---- Slide 4: AI Trends ----
        slide = add_dark_slide()
        add_textbox(slide, 0.8, 0.3, 11, 0.8, "AI/ML在血液病研究中的前沿趋势", 32, WHITE, True)
        add_textbox(slide, 0.8, 1.5, 11, 1.0, f"检索到AI相关论文: {len(ai_papers)}篇 (近3年)", 20, LIGHT_GRAY)
        ai_summary = "\n".join([f"  - {p.get('title', '')[:80]} ({p.get('year', '')})" for p in ai_papers[:5]])
        add_textbox(slide, 0.8, 2.5, 11, 4, f"代表性论文:\n{ai_summary}", 16, LIGHT_GRAY)

        # ---- Slide 5: Collaboration Topics ----
        slide = add_dark_slide()
        add_textbox(slide, 0.8, 0.3, 11, 0.8, "医工交叉合作课题方向", 32, WHITE, True)

        if topics:
            if isinstance(topics, list) and len(topics) > 0:
                if isinstance(topics[0], dict):
                    topic_lines = []
                    for i, t in enumerate(topics[:3], 1):
                        topic_lines.append(
                            f"{i}. {t.get('title', '')}\n"
                            f"   描述: {t.get('description', '')}\n"
                            f"   技术: {t.get('technologies', '')}\n"
                            f"   临床价值: {t.get('clinical_impact', '')}\n"
                            f"   可行性: {t.get('feasibility', '')}"
                        )
                    topics_text = "\n\n".join(topic_lines)
                else:
                    topics_text = "\n\n".join(topics[:3])
            add_textbox(slide, 0.8, 1.5, 11, 5, topics_text, 16, LIGHT_GRAY)
        else:
            add_textbox(slide, 0.8, 1.5, 11, 5, "暂无课题", 16, LIGHT_GRAY)

        # ---- Slide 6: Technical Roadmap ----
        slide = add_dark_slide()
        add_textbox(slide, 0.8, 0.3, 11, 0.8, "技术路线图", 32, WHITE, True)
        roadmap = (
            "阶段一 (0-6个月): 数据收集与预处理\n"
            "  - 建立临床数据库，收集多组学数据\n"
            "  - 数据标注与清洗\n\n"
            "阶段二 (6-12个月): 模型开发与验证\n"
            "  - 开发核心AI/ML算法\n"
            "  - 回顾性验证\n\n"
            "阶段三 (12-18个月): 临床应用与推广\n"
            "  - 前瞻性临床试验\n"
            "  - 产品化与注册"
        )
        add_textbox(slide, 0.8, 1.5, 11, 5, roadmap, 18, LIGHT_GRAY)

        # ---- Slide 7: Summary ----
        slide = add_dark_slide()
        add_textbox(slide, 1, 1, 11, 1.5, "总结", 36, WHITE, True, PP_ALIGN.CENTER)
        summary = (
            f"本建议基于对{expert_name}教授近5年论文的系统检索与分析，\n"
            f"结合血液病领域AI/ML前沿趋势，\n"
            f"提出{len(topics[:3])}个医工交叉合作课题方向。\n\n"
            "期待与专家深入交流，共同推进医工交叉研究！"
        )
        add_textbox(slide, 1, 3, 11, 3, summary, 20, LIGHT_GRAY, False, PP_ALIGN.CENTER)

        prs.save(ppt_path)
        print(f"  -> 已生成 {ppt_path}")
        return {"ppt_path": ppt_path, "current_step": "generate_ppt_completed"}

    except Exception as e:
        return {"error": f"生成PPT失败: {str(e)}"}